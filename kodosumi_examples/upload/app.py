import asyncio, os, re, time
from pathlib import Path
from tempfile import mkdtemp
from traceback import format_exc

import fastapi
import uvicorn
from pptx import Presentation
from ray import remote, serve

from kodosumi.core import ServeAPI, forms as F, Tracer, Launch
from kodosumi.response import Markdown

app = ServeAPI()

@remote
def process_file(file: str, tracer: Tracer, ignore_errors: bool = True):
    fs = tracer.fs_sync()
    tempfile = next(fs.download(file))
    tracer.debug_sync(f"start processing `{file}`")
    start_time = time.time()
    try:
        prs = Presentation(tempfile)
        text_runs = [
            run.text for slide in prs.slides 
            for shape in slide.shapes if shape.has_text_frame
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
        ]
        
        tracer.debug_sync(
            f"done processing `{file}` in {time.time() - start_time:.2f}s")
        
        tempdir = mkdtemp()
        md_file = Path(tempdir) / Path(file).with_suffix(".md").name
        content = re.sub(r"\s+", " ", " ".join(text_runs))
        
        with md_file.open("w") as f:
            f.write(content)
        
        fs.upload(str(md_file))
        os.remove(md_file)
        
        return {
            "source": file,
            "target": md_file.name,
            "length": len(text_runs),
            "words": len(content.split()),
            "runtime": time.time() - start_time,
            "error": None
        }
    except Exception as e:
        if not ignore_errors:
            raise RuntimeError(f"error processing `{file}`: {format_exc()}")
        return {
            "source": file,
            "target": None,
            "length": 0,
            "words": 0,
            "runtime": time.time() - start_time,
            "error": str(e)
        }
    finally:
        os.remove(tempfile)
        fs.close()

async def run(inputs: dict, tracer: Tracer):
    afs = await tracer.fs()
    files = await afs.ls("in")
    await tracer.markdown("### File Processing Started")
    
    file_links = [
        f"* [{f['path']}](/files/{tracer.fid}/{f['path']})" for f in files]
    await tracer.markdown("\n".join(file_links))
    
    futures = [
        process_file.remote(
            f['path'], tracer, ignore_errors=inputs['ignore_errors'])
        for f in files
    ]
    
    await afs.close()
    results = await asyncio.gather(*futures)
    
    output = ["### File Processing Completed"]
    for r in results:
        err = r['error']
        link = f"**ERROR:** {err}" if err else f"[markdown](/files/{tracer.fid}/out/{r['target']})"
        output.append(
            f"* {r['source']} with {r['length']} text runs in {r['runtime']:.2f}s - {link}"
        )
    
    return Markdown("\n".join(output))

model = F.Model(
    F.Markdown("""
    # Parallel Text Processing
    
    This application demonstrates the use of:
               
    * Forms with text input and file upload
    * `InputsFile` for file processing
    * Ray for parallel processing
    """),
    F.Break(),
    F.HR(),
    F.InputFiles(
        label="Upload Files", 
        name="files", 
        multiple=True, 
        directory=False, 
        required=False
    ),
    F.HR(),
    F.Submit("Start"),
    F.Cancel("Cancel"),
    F.Checkbox(
        name="ignore_errors",
        option="ignore errors",
        value=True
    ),
)

@app.enter(
    path="/", 
    model=model,
    summary="Text Processing with Ray",
    description="Demonstrates forms, InputsFile and Ray for parallel file processing.",
    version="0.1.0",
    author="example@kodosumi.com",
    organization="Kodosumi Examples",
    tags=["Test"]
)
async def enter(request: fastapi.Request, inputs: dict):
    return Launch(request, "kodosumi_examples.upload.app:run", inputs=inputs)

@serve.deployment
@serve.ingress(app)
class TextProcessor: pass

fast_app = TextProcessor.bind()  # type: ignore

if __name__ == "__main__":
    uvicorn.run(
        "kodosumi_examples.upload.app:app", 
        host="0.0.0.0", 
        port=8013, 
        reload=True
    ) 